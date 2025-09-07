from http.server import BaseHTTPRequestHandler
import json, urllib.parse, secrets, os, sqlite3
from http import cookies
from config import BASE_DIR, co, DB_PATH
from sessions import sessions
from chat_utils import get_prompt_by_mode, build_prompt
import base64
from docx import Document
import pdfplumber
import os

UPLOAD_DIR = os.path.join(BASE_DIR, "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)


# Conexión SQLite
def get_connection():
    return sqlite3.connect(DB_PATH)

# Función auxiliar para generar título con Cohere
def generate_title(text):
    try:
        prompt = f"Resume en 2 a 4 palabras el tema principal del siguiente mensaje: {text}\nTítulo:"
        response = co.generate(
            model='command-r-plus',
            prompt=prompt,
            max_tokens=10,
            temperature=0.5,
            k=0,
            p=0.75,
            frequency_penalty=0,
            presence_penalty=0,
            stop_sequences=["\n"],
        )
        return response.generations[0].text.strip()
    except Exception:
        return ''

class UnifiedHandler(BaseHTTPRequestHandler):
    def _send_cors_headers(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET,POST,OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')

    def _send_headers(self, status=200, content_type='text/html'):
        self.send_response(status)
        self.send_header('Content-type', content_type)
        self.end_headers()

    def do_OPTIONS(self):
        self.send_response(204)
        self._send_cors_headers()
        self.end_headers()

    def do_GET(self):
        if self.path == '/api/session':
            cookie_header = self.headers.get('Cookie')
            user_session = None
            if cookie_header:
                cookie = cookies.SimpleCookie()
                cookie.load(cookie_header)
                if 'session_id' in cookie:
                    session_id = cookie['session_id'].value
                    user_session = sessions.get(session_id)
            self._send_headers(200, 'application/json')
            self.wfile.write(json.dumps({'usuario': user_session}).encode('utf-8'))
            return

        if self.path == '/api/history':
            cookie_header = self.headers.get('Cookie')
            user_session = None
            if cookie_header:
                cookie = cookies.SimpleCookie()
                cookie.load(cookie_header)
                if 'session_id' in cookie:
                    session_id = cookie['session_id'].value
                    user_session = sessions.get(session_id)
            if not user_session:
                self._send_headers(401, 'application/json')
                self.wfile.write(json.dumps({'error': 'No autenticado'}).encode('utf-8'))
                return

            conn = get_connection()
            c = conn.cursor()
            c.execute(
                "SELECT id, mode, title, messages, created_at FROM chats WHERE username=? ORDER BY created_at ASC",
                (user_session,)
            )
            rows = c.fetchall()
            conn.close()

            chats = []
            for row in rows:
                chat_id, mode, title_db, messages_json, created_at = row
                try:
                    messages = json.loads(messages_json)
                except Exception:
                    messages = []
                title = title_db or ''
                chats.append({
                    'id': chat_id,
                    'mode': mode,
                    'title': title,
                    'msgs': messages,
                    'created_at': created_at
                })

            self._send_headers(200, 'application/json')
            self.wfile.write(json.dumps({'chats': chats}).encode('utf-8'))
            return

        # Servir archivos
        # Servir archivos
        if self.path in ['/', '/index.html']:
                filepath = os.path.join(BASE_DIR, 'index.html')
                content_type = 'text/html; charset=utf-8'
        elif self.path == '/registro.html':
                filepath = os.path.join(BASE_DIR, 'registro.html')
                content_type = 'text/html; charset=utf-8'
        elif self.path == '/menu.html':
                filepath = os.path.join(BASE_DIR, 'menu.html')
                content_type = 'text/html; charset=utf-8'
        elif self.path.endswith('.css'):
                filepath = os.path.join(BASE_DIR, self.path.lstrip('/'))
                content_type = 'text/css; charset=utf-8'
        elif self.path.endswith('.js'):
                filepath = os.path.join(BASE_DIR, self.path.lstrip('/'))
                content_type = 'application/javascript; charset=utf-8'
        elif self.path == "/telegram_webhook":
            content_length = int(self.headers.get('Content-Length', 0))
            post_data = self.rfile.read(content_length).decode("utf-8")
            from telegram_bot import process_update
            process_update(post_data)
            self.send_response(200)
            self.end_headers()
            return
        else:
            self.send_response(404)
            self.end_headers()
            return

        try:
            with open(filepath, 'rb') as f:
                content = f.read()
            self.send_response(200)
            self.send_header('Content-type', content_type)
            self.end_headers()
            self.wfile.write(content)
        except Exception:
            self.send_response(500)
            self.end_headers()

    def do_POST(self):
        content_length = int(self.headers.get('Content-Length', 0))
        post_data = self.rfile.read(content_length).decode('utf-8')

        # API Chat
        if self.path == '/api/chat':
                data = json.loads(post_data)
                messages = data.get('messages', [])
                mode = data.get('mode', 'general')
                chat_id = data.get('chat_id')
                title = data.get('title', '')

                # 📂 Procesar archivos enviados (Word, PDF, Excel, PowerPoint y TXT)
                for msg in messages:
                    if msg.get("sender") == "user" and msg.get("files"):
                        for f in msg["files"]:
                            filename = f.get("name")
                            filedata = f.get("data")
                            if filename and filedata:
                                file_bytes = base64.b64decode(filedata.split(",", 1)[-1])
                                filepath = os.path.join(UPLOAD_DIR, filename)
                                with open(filepath, "wb") as out_file:
                                    out_file.write(file_bytes)
                                print(f"📂 Archivo guardado: {filepath}")

                                extracted_text = ""

                                # PDF
                                if filename.lower().endswith(".pdf"):
                                    try:
                                        import pdfplumber
                                        with pdfplumber.open(filepath) as pdf:
                                            for page in pdf.pages:
                                                extracted_text += (page.extract_text() or "") + "\n"
                                    except Exception as e:
                                        print(f"⚠️ Error leyendo PDF: {e}")

                                # Word
                                elif filename.lower().endswith(".docx"):
                                    try:
                                        from docx import Document
                                        doc = Document(filepath)
                                        extracted_text = "\n".join(p.text for p in doc.paragraphs)
                                    except Exception as e:
                                        print(f"⚠️ Error leyendo DOCX: {e}")

                                # Excel
                                elif filename.lower().endswith((".xlsx", ".xls")):
                                    try:
                                        import openpyxl
                                        wb = openpyxl.load_workbook(filepath, data_only=True)
                                        sheet = wb.active
                                        for row in sheet.iter_rows(values_only=True):
                                            extracted_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
                                    except Exception as e:
                                        print(f"⚠️ Error leyendo Excel: {e}")

                                # PowerPoint
                                elif filename.lower().endswith(".pptx"):
                                    try:
                                        from pptx import Presentation
                                        prs = Presentation(filepath)
                                        for slide in prs.slides:
                                            for shape in slide.shapes:
                                                if hasattr(shape, "text"):
                                                    extracted_text += shape.text + "\n"
                                    except Exception as e:
                                        print(f"⚠️ Error leyendo PowerPoint: {e}")

                                # TXT
                                elif filename.lower().endswith(".txt"):
                                    try:
                                        with open(filepath, "r", encoding="utf-8") as txt_file:
                                            extracted_text = txt_file.read()
                                    except Exception as e:
                                        print(f"⚠️ Error leyendo TXT: {e}")

                                # Agregar texto extraído al mensaje
                                if extracted_text.strip():
                                    msg["text"] = (msg.get("text", "") + "\n" + extracted_text).strip()

                # Construir prompt con el texto extraído ya añadido
                prompt = get_prompt_by_mode(mode)
                full_prompt = build_prompt(messages, prompt)
                try:
                    response = co.generate(
                        model='command-r-plus',
                        prompt=full_prompt,
                        max_tokens=350,
                        temperature=0.75,
                        k=0,
                        p=0.75,
                        frequency_penalty=0,
                        presence_penalty=0,
                        stop_sequences=["--"],
                    )
                    answer = response.generations[0].text.strip()
                except Exception:
                    answer = "Error al obtener respuesta de Cohere."

                messages.append({'text': answer, 'sender': 'bot'})

                cookie_header = self.headers.get('Cookie')
                user_session = None
                if cookie_header:
                    cookie = cookies.SimpleCookie()
                    cookie.load(cookie_header)
                    if 'session_id' in cookie:
                        session_id = cookie['session_id'].value
                        user_session = sessions.get(session_id)

                if user_session:
                    try:
                        conn = get_connection()
                        c = conn.cursor()
                        chat_json = json.dumps(messages, ensure_ascii=False)

                        # Nuevo chat: generar título si no hay
                        if not chat_id and not title:
                            first_user_msg = next((m for m in messages if m['sender']=='user'), None)
                            if first_user_msg:
                                title = generate_title(first_user_msg['text'])

                        if chat_id:
                            if not title:
                                c.execute("SELECT title FROM chats WHERE id=? AND username=?", (chat_id, user_session))
                                row = c.fetchone()
                                if row:
                                    title = row[0]
                            c.execute(
                                "UPDATE chats SET mode=?, messages=?, title=? WHERE id=? AND username=?",
                                (mode, chat_json, title, chat_id, user_session)
                            )
                        else:
                            c.execute(
                                "INSERT INTO chats (username, mode, title, messages) VALUES (?, ?, ?, ?)",
                                (user_session, mode, title, chat_json)
                            )
                            chat_id = c.lastrowid
                        conn.commit()
                        conn.close()
                    except Exception as e:
                        print(f"Error guardando chat en DB: {e}")

                self.send_response(200)
                self._send_cors_headers()
                self.send_header('Content-Type', 'application/json')
                self.end_headers()
                self.wfile.write(json.dumps({"message": answer, "chat_id": chat_id, "title": title}).encode('utf-8'))
                return
        # API Title
        elif self.path == '/api/title':
            data = json.loads(post_data)
            messages = data.get('messages', [])
            chat_id = data.get('chat_id')
            title = ''
            if messages:
                text = messages[0].get('text', '')
                title = generate_title(text)

            cookie_header = self.headers.get('Cookie')
            user_session = None
            if cookie_header:
                cookie = cookies.SimpleCookie()
                cookie.load(cookie_header)
                if 'session_id' in cookie:
                    session_id = cookie['session_id'].value
                    user_session = sessions.get(session_id)

            if user_session and chat_id:
                try:
                    conn = get_connection()
                    c = conn.cursor()
                    c.execute(
                        "UPDATE chats SET title=? WHERE id=? AND username=?",
                        (title, chat_id, user_session)
                    )
                    conn.commit()
                    conn.close()
                except Exception as e:
                    print(f"Error actualizando título en DB: {e}")

            self.send_response(200)
            self._send_cors_headers()
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({"title": title}).encode('utf-8'))
            return

        # Registro/Login
        elif self.path == '/':
            post_vars = urllib.parse.parse_qs(post_data)
            usuario = post_vars.get('usuario', [None])[0]
            password = post_vars.get('password', [None])[0]
            action = post_vars.get('action', [None])[0]

            if not usuario or not password or not action:
                self._send_headers(400, 'application/json')
                self.wfile.write(json.dumps({'error': 'Faltan campos obligatorios'}).encode())
                return

            conn = get_connection()
            c = conn.cursor()

            if action == 'login':
                c.execute("SELECT password FROM users WHERE username=?", (usuario,))
                row = c.fetchone()
                if row and row[0] == password:
                    session_id = secrets.token_hex(16)
                    sessions[session_id] = usuario

                    c.execute(
                        "SELECT id, mode, title, messages, created_at FROM chats WHERE username=? ORDER BY created_at ASC",
                        (usuario,)
                    )
                    rows = c.fetchall()
                    chats = []
                    for r in rows:
                        chat_id, mode, title_db, messages_json, created_at = r
                        try:
                            messages = json.loads(messages_json)
                        except Exception:
                            messages = []
                        title = title_db or ''
                        chats.append({
                            'id': chat_id,
                            'mode': mode,
                            'title': title,
                            'msgs': messages,
                            'created_at': created_at
                        })

                    self.send_response(200)
                    self._send_cors_headers()
                    self.send_header('Content-type', 'application/json')
                    self.send_header('Set-Cookie', f'session_id={session_id}; HttpOnly; Path=/')
                    self.end_headers()
                    self.wfile.write(json.dumps({
                        'success': True,
                        'message': 'Bienvenido',
                        'usuario': usuario,
                        'chats': chats
                    }).encode())
                else:
                    self._send_headers(200, 'application/json')
                    self.wfile.write(json.dumps({'error': 'Usuario no registrado o contraseña incorrecta'}).encode())

            elif action == 'signup':
                confirmar = post_vars.get('confirmar', [None])[0]
                if not confirmar:
                    self._send_headers(400, 'application/json')
                    self.wfile.write(json.dumps({'error': 'Debes confirmar la contraseña'}).encode())
                    conn.close()
                    return

                if password != confirmar:
                    self._send_headers(200, 'application/json')
                    self.wfile.write(json.dumps({'error': 'Las contraseñas no coinciden'}).encode())
                    conn.close()
                    return

                # Validación de contraseña
                if len(password) < 8:
                    self._send_headers(200, 'application/json')
                    self.wfile.write(json.dumps({'error': 'La contraseña debe tener al menos 8 caracteres'}).encode())
                    conn.close()
                    return

                if not any(c.isupper() for c in password):
                    self._send_headers(200, 'application/json')
                    self.wfile.write(json.dumps({'error': 'La contraseña debe contener al menos una letra mayúscula'}).encode())
                    conn.close()
                    return

                try:
                    c.execute(
                        "INSERT INTO users (username, password) VALUES (?, ?)",
                        (usuario, password)
                    )
                    conn.commit()
                    session_id = secrets.token_hex(16)
                    sessions[session_id] = usuario
                    self.send_response(200)
                    self._send_cors_headers()
                    self.send_header('Content-type', 'application/json')
                    self.send_header('Set-Cookie', f'session_id={session_id}; HttpOnly; Path=/')
                    self.end_headers()
                    self.wfile.write(json.dumps({'success': True, 'message': 'Usuario registrado con éxito', 'usuario': usuario}).encode())
                except Exception:
                    self._send_headers(200, 'application/json')
                    self.wfile.write(json.dumps({'error': 'Usuario ya registrado'}).encode())
            else:
                self._send_headers(400, 'application/json')
                self.wfile.write(json.dumps({'error': 'Acción desconocida'}).encode())

            conn.close()
            return
        else:
            self.send_response(404)
            self.end_headers()
